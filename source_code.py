import streamlit as st
import openai
import fitz  # PyMuPDF
from openpyxl import load_workbook
from docx import Document
from io import BytesIO
from openai import OpenAI
# from mistralai.client import MistralClient
# from langchain.prompts import PromptTemplate
from functools import partial
# import langdetect
import streamlit as st
import openai
from pptx import Presentation
import re
import vertexai
from vertexai.generative_models import GenerativeModel, GenerationConfig
from langdetect import detect
import langid
from functools import partial
from google.cloud import translate_v3
import tiktoken
import requests
import math
from google.oauth2 import service_account
#  from typing import List, Dict
import json
import tempfile
import os  # Make sure os is imported for tempfile path


# Configuration
st.set_page_config(layout="wide")

# --- Constants ---
OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]
# MISTRAL_API_KEY = st.secrets["MISTRAL_API_KEY"]
PASSWORD = st.secrets["MDM_PASSWORD"]
LOCATION = st.secrets["LOCATION"]

# --- Google Cloud Credentials Handling ---
gcp_service_account_info_str = None
temp_file_path = None

try:
    gcp_service_account_info_str = st.secrets["GOOGLE_APPLICATION_CREDENTIALS"]

    if not gcp_service_account_info_str:
        st.error("Google Cloud service account key (GOOGLE_APPLICATION_CREDENTIALS) is empty in secrets.toml.")
        st.stop()

    # Create a temporary file
    with tempfile.NamedTemporaryFile(mode='w+', delete=False, encoding='utf-8', suffix='.json') as temp_file:
        temp_file.write(gcp_service_account_info_str)
        temp_file_path = temp_file.name

    # Initialize credentials from the temporary file path
    # This is the primary and only way to load the credentials from the string in this method
    gcp_credentials = service_account.Credentials.from_service_account_file(temp_file_path)

    # Get project_id directly from the credentials object (more reliable)
    # The credentials object itself contains the project_id once loaded
    project_id = gcp_credentials.project_id
    if not project_id:
        # Fallback to parsing if project_id is not directly available, but it should be
        # This part should ideally not be reached if the credential file is valid
        st.warning("Project ID not found directly in credentials object, attempting parse from string.")
        gcp_service_account_info = json.loads(gcp_service_account_info_str)
        project_id = gcp_service_account_info.get("project_id")
        if not project_id:
            st.error("Project ID not found in Google Cloud service account key.")
            st.stop()


except KeyError:
    st.error("Google Cloud service account key (GOOGLE_APPLICATION_CREDENTIALS) not found in secrets.toml.")
    st.stop()
except Exception as e: # Catch any other unexpected errors during loading
    st.error(f"An unexpected error occurred during Google Cloud secret setup: {e}")
    st.stop()
finally:
    # IMPORTANT: Clean up the temporary file immediately after use
    if temp_file_path and os.path.exists(temp_file_path):
        os.remove(temp_file_path)


# --- Initialize other clients ---
client = OpenAI(api_key=OPENAI_API_KEY)
# mistral_client = MistralClient(api_key=MISTRAL_API_KEY)
vertexai.init(project=project_id, location=LOCATION, credentials=gcp_credentials)
translate_client = translate_v3.TranslationServiceClient(credentials=gcp_credentials)


def display_language_selection(key_suffix, is_source=False):
    """
    Affiche la sélection de langue avec un champ de saisie pour les cas incertains
    
    Args:
        key_suffix: Suffixe pour les clés Streamlit
        is_source: Si True, affiche un champ de saisie pour la langue source
    """
    # List of predefined language options
    languages = ["Dutch", "French", "English", "Other"]
    
    # Selectbox for choosing the language
    if is_source:
        language_choice = st.selectbox("Choose source language:", languages, index=1, key=f'source_lang_{key_suffix}')
    else:
        language_choice = st.selectbox("Choose target language:", languages, index=1, key=f'to_lang_{key_suffix}')
    
    # Conditional display of text input when "Other" is selected
    if language_choice == "Other":
        language = st.text_input("Please use ISO-code:")
       
    else:
        language = {"Dutch": "nl", "French": "fr", "English": "en"}.get(language_choice)
                
    return language
    
def detect_language(text):
    try:
        return langdetect.detect(text)
    except:
        return "Unable to detect language"

def detect_language(text):
    try:
        # Get ISO code and score
        lang_code, score = langid.classify(text)
        
        # Map ISO code → English language name
        language_map = {
            'nl': 'Dutch',
            'en': 'English',
            'fr': 'French',
            'de': 'German',
            'es': 'Spanish',
            'it': 'Italian',
            'pt': 'Portuguese',
            'ru': 'Russian',
            'ja': 'Japanese',
            'zh': 'Chinese',
            'ar': 'Arabic',
            'hi': 'Hindi',
            'ko': 'Korean',
            'tr': 'Turkish',
            'pl': 'Polish',
            'sv': 'Swedish',
            'da': 'Danish',
            'fi': 'Finnish',
            'no': 'Norwegian',
            'cs': 'Czech',
            'el': 'Greek',
            'hu': 'Hungarian',
        }
        
        # Get English name (fallback to code if unknown)
        lang_name = language_map.get(lang_code, lang_code)
        
        return lang_code, lang_name
    except Exception as e:
        # Fallback to French if detection fails
        return "fr", "French"


def display_temperature_slider(key_suffix):
    return st.slider('**Select a Temperature**', min_value=0.1, max_value=1.0, step=0.1, key=f'temp_{key_suffix}')

# Utility functions
def read_pdf(file):
    text = ''
    bytes_stream = BytesIO(file.read())
    with fitz.open(stream=bytes_stream, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    
    return text

def read_pptx(file):
    text = ''
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ' '
    return text

def read_excel(file):
    text = ''
    wb = load_workbook(filename=file)
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                text += str(cell) + ' '
    return text

def read_docx(file):
    doc = Document(file)
    return ' '.join([paragraph.text for paragraph in doc.paragraphs])


def display_file_uploader():
    MAX_FILE_SIZE_MB = 50  # 50MB limit
    
    uploaded_file = st.file_uploader(
        "Upload file (PDF, PPTX, XLSX, DOCX)", 
        type=['pdf', 'pptx', 'xlsx', 'docx'],
        accept_multiple_files=False,
        help=f"Maximum file size: {MAX_FILE_SIZE_MB}MB"
    )
    
    if uploaded_file:
        # Check file size
        if uploaded_file.size > MAX_FILE_SIZE_MB * 1024 * 1024:
            st.error(f"File too large. Maximum size is {MAX_FILE_SIZE_MB}MB")
            return ""
            
        try:
            if uploaded_file.type == "application/pdf":
                return read_pdf(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return read_pptx(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return read_excel(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return read_docx(uploaded_file)
        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            return ""
    return ""

def display_text_input(disp_key):
    return st.text_area('Or enter text to translate', height=150, key=disp_key)


def chunk_text(text: str, max_len: int = 5000) -> list[str]:
    """
    Splits a long text into smaller chunks based on sentence boundaries.

    Args:
        text: The input text.
        max_len: The maximum length (in characters) for a chunk.

    Returns:
        A list of text chunks.
    """
    sentences = re.split(r'(?<=[.?!])\s+(?=[A-Z])', text) # Splits on .?! followed by space and capital letter
    if not sentences:
        # Basic fallback: handle cases where splitting by sentence doesn't work
        # Split strictly by max_len if sentence splitting fails or text is shorter than one sentence
        st.write(f"INFO: Basic fallback chunking applied. Text length: {len(text)}, max_len: {max_len}")
        return [text[i:i+max_len] for i in range(0, len(text), max_len)] if text else []

    chunks = []
    current_chunk = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        # Check if the sentence itself is too long
        if len(sentence) > max_len:
            # If a sentence is too long, add the current chunk first (if any)
            if current_chunk:
                 chunks.append(current_chunk)
                 current_chunk = "" # Reset current chunk

            # Then, split the long sentence itself into chunks
            st.write(f"WARNING: Sentence longer than max_len ({len(sentence)} > {max_len}), splitting sentence.")
            for i in range(0, len(sentence), max_len):
                 chunks.append(sentence[i:i+max_len])
            # After splitting the long sentence, there's no current_chunk to continue with
            continue # Move to the next sentence

        # Check if adding the new sentence would exceed max_len
        if len(current_chunk) + len(sentence) + (1 if current_chunk else 0) <= max_len:
            if current_chunk:
                current_chunk += " " + sentence
            else:
                current_chunk = sentence
        else:
            # Current chunk is full, add it to the list and start a new chunk
            if current_chunk: # Prevent empty chunks
                 chunks.append(current_chunk)
            current_chunk = sentence

    # Add the last chunk if it's not empty
    if current_chunk:
        chunks.append(current_chunk)

    # Fallback if chunking somehow resulted in empty list but text exists
    if not chunks and text:
         st.write(f"WARNING: Sentence splitting resulted in no chunks, falling back to hard chunking. Text length: {len(text)}, max_len: {max_len}")
         return [text[i:i+max_len] for i in range(0, len(text), max_len)]
    elif not chunks and not text:
         return [] # Empty text should result in empty chunks

    return chunks


def translate_text_with_v3(
    text: str = "YOUR_TEXT_TO_TRANSLATE",
    source_language_code: str = "nl",
    target_language_code: str = "fr",
    max_chunk_len: int = 5000, # Max length for chunks (karakters)
    project_id: str = project_id, # This project_id should also come from the global scope
    translate_client_obj: translate_v3.TranslationServiceClient = None 
) -> str:
    """Translates text using chunking, suitable for long inputs.

    Args:
        text: The content to translate.
        source_language_code: The code of the source language.
        target_language_code: The code of the target language.
        max_chunk_len: Maximum character length for each chunk. Google Cloud
                        Translate v3 has limits (e.g., 30k *bytes* per request),
                        so keep this reasonably low (e.g., 5000-15000 characters).
        project_id: Your Google Cloud Project ID.
        translate_client_obj: The initialized Google Cloud TranslationServiceClient object.

    Returns:
        The translated text as a single string.

    Raises:
        ValueError: If project_id is not provided and cannot be found in environment.
        Exception: If the translation API call fails.
    """
    if not text:
        return ""

    # Ensure the client object is provided
    if translate_client_obj is None:
        raise ValueError("TranslationServiceClient object must be provided to translate_text_with_v3.")

    # Chunk de tekst
    text_chunks = chunk_text(text, max_len=max_chunk_len)

    # Use the passed client object
    client = translate_client_obj
    parent = f"projects/{project_id}/locations/global"
    mime_type = "text/plain"

    translated_pieces = []
    for chunk in text_chunks:
        response = client.translate_text(
            contents=[chunk],
            parent=parent,
            mime_type=mime_type,
            source_language_code=source_language_code,
            target_language_code=target_language_code,
        )
        # response.translations is een lijst met één element
        translated_pieces.append(response.translations[0].translated_text)

    return " ".join(translated_pieces)


def create_word_document(
    text1: str, # The refined text 
    text2: str         # The original text 
) -> bytes: # Always returns bytes if text2 has content, handles empty/None text1
    """
    Creates a Word document (.docx) containing two distinct text sections
    with introductory headers. Handles simple paragraph and line breaks within the texts.

    If text1 (the refined translation) is empty or None, a placeholder text is inserted.
    Assumes text2 (the basic translation) will always have content.

    Args:
        text1: The content for the first section (e.g., the refined translation). Can be empty or None.
        text2: The content for the second section (e.g., the basic translation).

    Returns:
        Bytes of the generated Word document. Returns None only if an error occurs during document creation
        (though typically errors should raise exceptions).
    """
    # Based on user's clarification, we don't check if text1 is missing to return None.
    # We assume text2 has content, so we proceed with document creation.

    try:
        document = Document()

        # --- Add the first section (Refined Translation) ---
        document.add_paragraph("Hereafter the latest refined translation:")
        document.add_paragraph() # Add an empty line for separation

        # Check if text1 is empty or None, insert placeholder if so
        if not text1: # This handles both None and empty string
            document.add_paragraph("There is no refined translation.")
        else:
            # Add the content of text1 if it exists, applying paragraph/line break logic
            paragraphs1 = text1.split('\n\n')
            for para_text1 in paragraphs1:
                if para_text1.strip():
                    p1 = document.add_paragraph()
                    lines1 = para_text1.split('\n')
                    for i, line1 in enumerate(lines1):
                        run1 = p1.add_run(line1)
                        if i < len(lines1) - 1:
                            run1.add_break() # Add a 'soft' line break

        # Add a blank line or separator after the first section content
        document.add_paragraph() # Add an empty line

        # --- Add the second section (Basic Translation) ---
        document.add_paragraph("Hereafter the latest translation without refinement:")
        document.add_paragraph() # Add an empty line for separation

        if not text2: # This handles both None and empty string
            document.add_paragraph("There is no original translation.")
        
        # Add the content of text2, applying paragraph/line break logic
        # We assume text2 has content based on user's clarification, so we proceed directly
        paragraphs2 = text2.split('\n\n')
        for para_text2 in paragraphs2:
            if para_text2.strip():
                p2 = document.add_paragraph()
                lines2 = para_text2.split('\n')
                for i, line2 in enumerate(lines2):
                    run2 = p2.add_run(line2)
                    if i < len(lines2) - 1:
                        run2.add_break() # Add a 'soft' line break
                                 
        # --- Save the document to a buffer ---
        buffer = BytesIO()
        document.save(buffer)
        buffer.seek(0) # Go to the beginning of the buffer
        return buffer.getvalue() # Return the bytes

    except Exception as e:
        # Handle potential errors during document creation
        st.write(f"Error creating Word document: {e}") 
        return None # Return None on error during processing


def translate_text_vertexai_prompting(
    text: str,
    source_language_code: str,
    target_language_code: str,
    instructions: str,
    temperature: float,
    model_name: str = "gemini-2.5-pro-preview-05-06", 
    max_chunk_len: int = 200000, # Adjust based on model's context window minus prompt length
) -> str:
    """
    Translates text using a Vertex AI LLM with custom prompting.

    Args:
        text: The text to be translated.
        source_language_code: Source language code (e.g., "en").
        target_language_code: Target language code (e.g., "fr").
        prompt_instructions: Specific instructions for the LLM (e.g., tone, style, glossary).
        model_name: The Vertex AI model name to use .
        max_chunk_len: Maximum length of a text chunk for the LLM's context window.
                       Must be less than the model's total context window minus prompt size.
        temperature: Model's temperature for output randomness (0.0 to 1.0).

    Returns:
        The translated text as a string.
        Returns an error message if the translation fails.
    """
    if not text:
        return "" # No text to translate

    # Load the chosen model
    try:
        # Use GenerativeModel for Gemini models
        model = GenerativeModel(model_name)
        
    except Exception as e:
        st.write(f"ERROR: Failed to load Vertex AI model '{model_name}': {e}")
        return f"[TRANSLATION FAILED - MODEL LOADING ERROR: {model_name}]"

       
    text_chunks = chunk_text(text, max_chunk_len)

    if not text_chunks:
         st.write("ERROR: Could not split text into chunks")
         return "[TRANSLATION FAILED - CHUNKING ERROR]"

    translated_pieces = []

    # --- Translate each chunk ---
    for i, chunk in enumerate(text_chunks):

        if source_language_code == 'nl': 
            full_prompt = f"""
            Je bent een professionele senior vertaler van het Nederlands naar {target_language_code}.
            Vertaal de onderstaande tekst zorgvuldig. Hou rekening met de eventuele instructies.
            Zorg ervoor dat de vertaling natuurlijk, vlot, helder en engagerend is.  
    
            Tekst:
            ---
            {chunk}
            ---

            Instructies:
            ---
            {instructions}
            ---
                       
            Herlees je eerste vertaling en stel jezelf de volgende vragen: 
            - Helderheid en directheid: Zijn de zinnen duidelijk en to the point? Kan de boodschap helderder en directer worden verwoord? 
            - Natuurlijke flow: Sluiten de zinnen en alinea's vlot aan? Zijn er haperingen in de vertaling?
            - Precise woordkeuze: Zijn de gekozen woorden de meest passende gezien de context en kernboodschap van de tekst?
            - Consistentie: Wordt dezelfde terminologie gebruikt en is de toon door het hele artikel gelijk?
            - Structuur: Is de informatie logisch opgebouwd en gemakkelijk begrijpbaar?
            
            Verbeter de eerste vertaling in functie van je antwoorden op deze kritische vragen.
            Kijk ook na of de instructies nog wel kloppen. 

            Pas de lay-out aan zodat de vertaalde inhoud duidelijk en effectief wordt weergegeven. Behoud de structuur en opmaak van de brontekst (zoals alinea’s, bullets, koppen, inline-opmaak in HTML/Markdown en codeblokken) als deze van hoge kwaliteit zijn en bijdragen aan de Nederlandse versie. Is dit niet het geval, of indien aanpassing nodig is, optimaliseer dan de lay-out voor een maximale leesbaarheid en een heldere presentatie in het Nederlands.
            
            Geef de finaal verbeterde versie door als antwoord, zonder extra tekst of commentaar.
            
            Vertaling:
            """
            
        elif source_language_code == 'fr':
            full_prompt = f"""
            Tu es un traducteur professionnel de premier ordre du français vers le {target_language_code}. 
            Traduis le texte ci-dessous avec le plus grand soin, en tenant compte des instructions. 
            Assure-toi que la traduction soit naturelle, fluide, limpide et engageante. 
            
            Texte :
            ---
            {chunk}
            ---

            Instructions:
            ---
            {instructions}
            ---
            
            Relis ta première traduction et pose-toi les questions suivantes:
            - Clarté et concision : Les phrases sont-elles claires et bien conçues ? Le message peut-il être exprimé de façon plus efficace ? 
            - Fluidité naturelle : Les idées s'enchaînent-elles naturellement ? La traduction est-elle quelque part hésitante ou cahotante?
            - Précision du vocabulaire : a-t-on utilisé les bons mots, les expressions appropriées, les bonnes tournures au vu du contexte et des messages-clé?
            - Cohérence : La terminologie et le ton sont-ils cohérents du début à la fin ?
            - Structure : L'information est-elle logiquement structurée et facile à comprendre ?
            
            Améliore sensiblement la qualité de la première traduction selon les réponses à ces questions.
            Assure-toi que le texte respecte toujours les instructions.

            Consignes de mise en page pour la traduction:

            Adapte la mise en page afin de présenter le contenu traduit de manière claire et efficace. Si la mise en forme du texte source (tels que les sauts de paragraphe, listes à puces, titres, balisage HTML/Markdown ou extraits de code) est de qualité et améliore la compréhension en français, conserve-la. Dans le cas contraire, ou si des ajustements sont nécessaires, privilégie une présentation guidée par le contenu pour garantir une lisibilité et une clarté optimales en français.
            
            Donne uniquement cette version finale, sans autre texte ni commentaire.
            
            Traduction :
            """

        else:
            full_prompt = f"""
            You are a senior professional translator from {source_language_code} to {target_language_code}.
            Translate the following text with care, ensuring your translation is natural, fluent, clear, and engaging.

            Text :
            ---
            {chunk}
            ---

            Instructions:
            ---
            {instructions}
            ---
            
            After completing your initial translation, review your work and consider the following questions:
            
            - Are the sentences clear, concise, and direct? Can any ideas be expressed more simply and with greater impact?
            - Do the sentences and paragraphs flow naturally?
            - Have you chosen the most appropriate words and terminologies?
            - Is terminology and tone consistent throughout?
            - Is the information logically structured and easy to follow?
            
            Revise your translation accordingly.
            
            Adapt the layout to effectively present the translated content. If the source text's formatting (paragraph breaks, bullet lists, headers, inline markup (HTML/Markdown), and code snippets) is of high quality and enhances the {target_language_code} version, preserve it. Otherwise, or in addition, ensure the layout is guided by the content for optimal readability and presentation in {target_language_code}.

            Please provide only the final, improved version—no extra commentary is needed.
            
            Translation:
            """
            
        try:
            # For Gemini models (using GenerativeModel)
            response = model.generate_content(
                full_prompt,
                generation_config={
                    "temperature": temperature
                }
            )

            # Extract the translated text from the response
            translated_chunk = ""
            if response and hasattr(response, 'text'): # For Gemini
                 translated_chunk = response.text.strip()
            else:
                 st.write(f"WARNING: Unexpected response structure for chunk {i+1}. Response: {response}")


            if translated_chunk:
                 translated_pieces.append(translated_chunk)
                
            else:
                 st.write(f"WARNING: Received empty translation for chunk {i+1}. Full response: {response}")
                 translated_pieces.append("[CHUNK TRANSLATION FAILED - EMPTY]")


        except Exception as e:
            st.write(f"ERROR translating chunk {i+1} with Vertex AI: {e}")
            translated_pieces.append(f"[CHUNK {i+1} TRANSLATION ERROR: {e}]")
            return f"[TRANSLATION FAILED ENTIRELY AT CHUNK {i+1}]"

    # --- Join results ---
    full_translation = " ".join(translated_pieces)
    return full_translation

def get_openai_translation(source_text, source_language, target_language, temp_choice, model="gpt-4.1"):

    instruction = f"""
    You are a master translator with native-level fluency in both {source_language} and {target_language}, plus a deep grasp of each culture’s idioms, tone, and stylistic nuances. 
    Translate any text from {source_language} into fluid, engaging {target_language} that reads as though it were originally crafted in {target_language} by a seasoned writer—faithfully conveying the author’s intent while adapting references and phrasing so naturally that no trace of translation remains.
    """

    text_chunks = chunk_text(source_text, max_len=120000)

    
    translated_pieces = []
    for chunk in text_chunks:  
        prompt = f"""
        Translate the text between the triple chevrons from {source_language} into {target_language}.
    
        <<<
        {chunk}
        >>>
        
        Guidelines
        1. **Meaning & Intent** Convey the author’s message faithfully and convincingly.
        2. **Tone & Register** Adopt the tone of the original.
        3. **Cultural Adaptation** Render idioms, metaphors, and cultural references naturally for {target_language} readers. If no equivalent exists, use a culturally apt substitute.
        4. **Terminology** Use field-specific terms accurately and consistently. 
        5. **Proper Nouns & Trademarks** Keep them in the original language unless an established localized form exists.
        6. **Formatting & Layout**: Adapt the layout to effectively present the translated content. If the source text's formatting (paragraph breaks, bullet lists, headers, inline markup (HTML/Markdown), and code snippets) is of high quality and enhances the {target_language} version, preserve it. Otherwise, or in addition, ensure the layout is guided by the content for optimal readability and presentation in {target_language}.
        7. **Ambiguity** If the source is genuinely ambiguous, retain that ambiguity; only insert a brief [translator note] when clarification is essential.
        8. **Clarity & Fluency** The result should read as if originally written in {target_language} by a skilled writer-editor.
        
        **Output**  
        Return **only** the translated text—no comments.
        """
        
        response = client.responses.create(
            model=model,
            # previous_response_id=response.id,
            instructions=instruction,
            input=prompt,
            temperature=temp_choice
        )
        
        translated_pieces.append(response.output_text)

    return " ".join(translated_pieces)

def count_tokens(text: str, model: str = "gpt-4o") -> int:
    enc = tiktoken.encoding_for_model(model)
    return len(enc.encode(text))


def get_openai_refinement(version_brute, version_travaillee, target_language, gloss_instructions, model="o4-mini"):

    if count_tokens(version_brute) >= 90000:
        return "Your text is too long. Please split it or use GEMINI."
    
    instruction = f"""
    You are an exceptional {target_language} writer-editor. Produce clear, fluent, engaging text by implementing concise, actionable improvements to given translations.
    """

    prompt = f"""
    # INPUT
    ## Draft A — Raw translation
    <<<
    {version_brute}
    >>>
    
    ## Draft B — Edited translation
    <<<
    {version_travaillee}
    >>>
    
    # TASK
    Create one definitive version in {target_language}.  
    If A and B differ, default to Draft B unless Draft A offers a clearer or more accurate wording.
    
    # QUALITY CRITERIA
    1. Engagement — varied rhythm and vivid, natural word choice that captures attention.  
    2. Fluency — seamless paragraph flow; no awkward transitions or repetitions.  
    3. Clarity — exact meaning, precise vocabulary; add **no** new ideas.  
    4. Consistent register — uniform tone from start to finish.  
    5. Domain accuracy — apply glossary terms.
    
    # GLOSSARY & SPECIAL INSTRUCTIONS 
    {gloss_instructions}
    
    # OUTPUT FORMAT
    - Keep the logical hierarchy (headings → sub-headings → paragraphs → lists).  
    - Improve layout (spacing, bullet marks) to maximise readability.  
    - Remove extraction artefacts (garbled strings, header/footer debris, missing spaces).
    - Return **only** the final text—no comments, tags, or explanations.
    
    DO NOT OUTPUT ANYTHING EXCEPT THE FINAL REVISED TEXT .
    """      
        
    response = client.responses.create(
        model=model,
        # previous_response_id=response.id,
        instructions=instruction,
        input=prompt
    )
    
    return response.output_text

def improve_translation_with_gpt(original_text, translated_text, guidelines, target_lang, term_instructions): 

    instruction = f"You are a senior translator-editor in {target_lang}. Your mission: craft publication-ready prose that is clear, fluent, engaging, and faithful to provided guidelines."
           
    prompt = f"""
    Below, you will find the original text, an initial translation, and potentially user feedback and glossary guidelines for improvement.

    ---
    Original Text:
    {original_text}
    ---

    Initial Translation:
    {translated_text}
    
    ---
    User Feedback:
    {guidelines}

    ---
    Glossary Guidelines:
    {term_instructions}

    
    **Task:** 
    Generate a final, improved translated text in {target_lang}.
    
    **Overall Principle:** 
    Do not feel limited by the perceived quality or style of the Original Text. Strive to create a text of the highest intrinsic quality, as if originally written in the target language, even if this means the translation reads better, is more effective, or uses culturally more appropriate phrasing/idioms than the original source text. Aim for engagement, fluency and clarity in the target language.
    
    **Instructions:**
    1.  Read and thoroughly understand the Original Text and the Initial Translation.
    2.  **IF the "User Feedback" section contains specific instructions or suggestions:**
        - Carefully analyze the feedback and its intent.
        - Apply the analyzed feedback to the Initial Translation.
    3.  **IF the "User Feedback" section is empty or contains no specific instructions (e.g., just filler text like "looks good" or is left blank):**
        - Ignore the "User Feedback" section.
    4.  In either case, follow these guidelines:
        - Improve fluency, engagement, clarity, rhythm and tone of the translation.
        - Adhere to the Overall Principle of striving for the highest intrinsic quality.
        - When feedback and glossary clash, glossary wins. {term_instructions}
            
    **Output Format:**
    DO NOT OUTPUT ANYTHING EXCEPT THE FINAL REVISED TEXT.
    """
        
    response = client.responses.create(
        model="o3",
        # previous_response_id=response.id,
        instructions=instruction,
        input=prompt
    )
    
    return response.output_text

    
def improve_translation_with_gemini(original_text, translated_text, guidelines, term_instructions, target_lang):
    
    # Load the chosen model
    try:
        # Use GenerativeModel for Gemini models
        model = GenerativeModel("gemini-2.5-pro-preview-05-06")
        
    except Exception as e:
        st.write(f"ERROR: Failed to load Vertex AI model: {e}")
        return f"[TRANSLATION FAILED - MODEL LOADING ERROR]"

    full_prompt = f"""
    You are a senior translator-editor in {target_lang}. Your mission: craft publication-ready prose that is clear, fluent, engaging, and faithful to provided guidelines.
    Below, you will find the original text, a proposed translation, and potentially user feedback and glossary guidelines for improvement.

    ---
    Original Text:
    {original_text}
    ---

    Translation:
    {translated_text}
    
    ---
    User Feedback:
    {guidelines}

    ---
    Glossary Guidelines:
    {term_instructions}

    
    **Task:** 
    Generate a final, improved translated text in {target_lang}.
    
    **Overall Principle:** 
    Do not feel limited by the perceived quality or style of the Original Text. Strive to create a text of the highest intrinsic quality, as if originally written in the target language, even if this means the translation reads better, is more effective, or uses culturally more appropriate phrasing/idioms than the original source text. Aim for engagement, fluency and clarity in the target language.
    
    **Instructions:**
    1.  Read and thoroughly understand the Original Text and the Initial Translation.
    2.  **IF the "User Feedback" section contains specific instructions or suggestions:**
        - Carefully analyze the feedback and its intent.
        - Apply the analyzed feedback to the Initial Translation.
    3.  **IF the "User Feedback" section is empty or contains no specific instructions (e.g., just filler text like "looks good" or is left blank):**
        - Ignore the "User Feedback" section.
    4.  In either case, follow these guidelines:
        - Improve fluency, engagement, clarity, rhythm and tone of the translation.
        - Adhere to the Overall Principle of striving for the highest intrinsic quality.
        - When feedback and glossary clash, glossary wins. {term_instructions}
            
    **Output Format:**
    DO NOT OUTPUT ANYTHING EXCEPT THE FINAL REVISED TEXT.
    """
        
    response = model.generate_content(
                full_prompt,
                generation_config={
                    "temperature": 0.6
                }
            )
    return response.text.strip()


def check_glossary(
    source_text: str,
    source_language: str,
    target_language: str,
    glossary_data: list, 
    temperature: float = 0.1,
    model_name: str = "gemini-2.0-flash"   
) -> dict:
    
    """
    Gebruikt een Vertex AI LLM om termen uit een brontekst te extraheren die overeenkomen
    met termen in een opgegeven glossarium, en geeft de aanbevolen vertaling.
    Als de source_language geen overeenkomstige termen in het glossarium heeft,
    of als voor een gevonden term geen vertaling naar target_language bestaat,
    wordt respectievelijk een lege lijst van matches of een lege string voor de vertaling gegeven.

    Args:
        project_id: Je Google Cloud Project ID.
        model_name: De naam van het te gebruiken LLM (bv. "gemini-1.5-flash-001").
        source_text: De brontekst die geanalyseerd moet worden.
        source_language: De taalcode van de brontekst (bv. "fr", "nl", "en").
                         De effectiviteit hangt af van de aanwezigheid van deze taalcode
                         als key in de entries van glossary_data.
        target_language: De taalcode voor de aanbevolen vertaling (bv. "fr", "nl", "en").
                         De effectiviteit hangt af van de aanwezigheid van deze taalcode
                         als key in de entries van glossary_data.
        glossary_data: Een Python-lijst van dictionaries, waarbij elke dictionary
                       een term in verschillende talen bevat. Bijvoorbeeld:
                       [{"fr": "Terme FR", "nl": "Term NL", "en": "Term EN"}, ...].
        temperature: De sampling temperatuur voor de LLM.
        location: De Google Cloud regio (bv. "us-central1").

    Returns:
        Een Python dictionary.
        Bij succes (ook als er geen matches zijn) bevat het de volgende structuur:
        {
          "matches": [
            {
              "source_text_term": "De term zoals letterlijk gevonden in de brontekst",
              "corresponding_glossary_term_source_lang": "De corresponderende glossariumterm in de brontaal",
              "preferred_translation_target_lang": "De aanbevolen vertaling (lege string "" indien niet beschikbaar)"
            },
            // ... meer matches ...
          ]
        }
        Bij een kritieke fout (bv. configuratie, LLM-fout): {"error": "Foutmelding", "raw_output": "..."}
    """
    
    try:
        model = GenerativeModel(model_name)

        # Stel de lijst met glossariumtermen voor de prompt samen op basis van source_language
        glossary_terms_for_prompt = []
        for entry in glossary_data:
            if isinstance(entry, dict) and source_language in entry and entry[source_language]:
                glossary_terms_for_prompt.append(str(entry[source_language]))

        # Als er geen glossariumtermen zijn voor de opgegeven brontaal,
        # dan kunnen we geen matches vinden. Sla LLM call over en retourneer lege matches.
        if not glossary_terms_for_prompt:
            return {"matches": []}

        glossary_list_for_prompt = "- " + "\n- ".join(glossary_terms_for_prompt)
        
        
        prompt = f"""
        You are an AI assistant specialized in text analysis and terminology recognition.
        Your task is to analyze the following source text (in '{source_language}') and identify all segments that are identical to, or very similar to (e.g., singular/plural, minor spelling errors, closely related synonyms, or variations in word order for multi-word terms) the terms in the glossary provided below.

        The glossary below contains terms in '{source_language}'.
        Glossary:
        {glossary_list_for_prompt}

        Source Text (in '{source_language}'):
        ---
        {source_text}
        ---

        Output Instructions:
        1. The output must be a valid JSON object.
        2. This JSON object must contain a single key: "matches".
        3. The value of "matches" must be a list of JSON objects.
        4. Each JSON object in the "matches" list represents a found correspondence and must contain the following three string keys:
           - "source_text_term": The exact segment of text as it appears literally in the source text that matches a glossary term.
           - "corresponding_glossary_term_source_lang": The specific term from the provided glossary (in '{source_language}') that matches the "source_text_term".
           - "preferred_translation_target_lang": You will determine this value. Look up the "corresponding_glossary_term_source_lang" in the full glossary data structure provided below. Find its translation into the target language '{target_language}'. If a translation for that specific term into '{target_language}' exists in the full glossary, provide it. If it does not exist, or if the term itself is missing for that language, use an empty string "" for this value.
        5. If no matching terms are found in the source text, the value of "matches" must be an empty list, i.e., {{"matches": []}}.
        6. Ensure the output is **only** the JSON object, without any extra text, explanation, or markdown formatting (like ```json ... ```) before or after it.

        Full glossary data structure (This is for your reference to find the translations into '{target_language}'. Do NOT use this structure for matching against the source text directly; use the "Glossary" list above for that.):
        {json.dumps(glossary_data, indent=2, ensure_ascii=False)}

        Begin the JSON output here:
        """

        generation_config = GenerationConfig(
            temperature=temperature,
            # response_mime_type="application/json"
        )

        response = model.generate_content(
            prompt,
            generation_config=generation_config
        )

        cleaned_response_text = response.text.strip().removeprefix("```json").removesuffix("```").strip()
        
        if not cleaned_response_text:
            # Hier kiezen we voor geen matches.
            return {"matches": []}

        parsed_json = json.loads(cleaned_response_text)

        if not isinstance(parsed_json, dict) or "matches" not in parsed_json:
            st.write(f"Warning: LLM output is not in expected JSON format (missing 'matches' key). Raw: {cleaned_response_text}")
            return {"matches": []}
        if not isinstance(parsed_json["matches"], list):
            st.write(f"Warning: LLM output 'matches' key is not a list. Raw: {cleaned_response_text}")
            return {"matches": []}

        validated_matches = []
        for match in parsed_json.get("matches", []):
            if isinstance(match, dict) and \
               "source_text_term" in match and \
               "corresponding_glossary_term_source_lang" in match:
                
                source_text_term_llm = str(match["source_text_term"])
                glossary_term_src_lang_llm = str(match["corresponding_glossary_term_source_lang"])
                
                actual_preferred_translation = ""
                # Zoek de vertaling op in de originele glossary_data
                for g_entry in glossary_data:
                    if isinstance(g_entry, dict) and g_entry.get(source_language) == glossary_term_src_lang_llm:
                        actual_preferred_translation = g_entry.get(target_language, "") # Geeft "" als target_language key niet bestaat of waarde None is
                        break
                
                validated_matches.append({
                    "source_text_term": source_text_term_llm,
                    "corresponding_glossary_term_source_lang": glossary_term_src_lang_llm,
                    "preferred_translation_target_lang": actual_preferred_translation
                })
            else:
                st.write(f"Warning: Invalid match object is dismissed: {match}")

        return {"matches": validated_matches}

    except json.JSONDecodeError as e:
        # Als JSON parsen faalt, is de output van LLM waarschijnlijk corrupt.
        raw_output = cleaned_response_text if 'cleaned_response_text' in locals() else 'Response text not available'
        st.write(f"Critical error: could not parse LLM JSON output: {e}. Raw output: {raw_output}")
        return {"error": f"Failed to parse LLM JSON output: {e}", "raw_output": raw_output}
    except ImportError as e: # Mocht vertexai of andere imports falen op een hoger niveau
        st.write(f"Critical error: Missing import: {e}")
        return {"error": f"Missing import: {e}"}
    except Exception as e:
        # Vang andere onverwachte Vertex AI / model errors of andere exceptions
        # Dit zijn waarschijnlijk meer kritieke fouten dan "geen match gevonden".
        import traceback
        st.write(f"Critical Error: An unexpected error occurred: {e}\n{traceback.format_exc()}")
        return {"error": f"An unexpected error occurred: {e}"}

def format_terminology_for_prompt(glossary_analysis: dict, source_lang_name: str, target_lang_name: str) -> str:
    """
    Formats the output of check_glossary_v4 into a string
    that can be used in the prompt of another agent.
    """
    if "error" in glossary_analysis:
        return f"Note: There was an error checking terminology ({glossary_analysis['error']}). No specific terminology guidelines are available."

    matches = glossary_analysis.get("matches", [])
    if not matches:
        return "Note: No specific terms from the glossary were found in the source text, or no specific terminology guidelines are available."

    instructions = [f"When processing the text, take into account the following recommended terminology (from {source_lang_name} to {target_lang_name}):"]

    has_specific_translations = False
    for match in matches:
        # The term as it might appear in the source text
        source_term_context = match['source_text_term']
        # The official glossary term in the source language
        glossary_term_source = match['corresponding_glossary_term_source_lang']
        # The recommended translation in the target language
        preferred_translation = match['preferred_translation_target_lang']

        if preferred_translation: # Only if there is an actual preferred translation
            instructions.append(
                f"- For the {source_lang_name} term '{glossary_term_source}' "
                f"(potentially appearing as '{source_term_context}' in the text), "
                f"use the {target_lang_name} translation: '{preferred_translation}'."
                "If multiple options are available for a specific term, choose the one that seems most suitable."
            )
            has_specific_translations = True
        elif glossary_term_source : # A source language term was matched, but no target language translation
            instructions.append(
                f"- The {source_lang_name} term '{glossary_term_source}' "
                f"(potentially appearing as '{source_term_context}' in the text) has been recognized. "
                f"No specific preferred translation to {target_lang_name} is provided in the glossary."
            )

    if not has_specific_translations and any(match['corresponding_glossary_term_source_lang'] for match in matches):
        # There were matches, but none had a concrete translation.
        # instructions.append("\nAlthough some terms were recognized, no specific preferred translations were provided.")
        pass # The individual messages above already cover this.
    elif not has_specific_translations:
        return "Note: No specific terms from the glossary with a preferred translation were found in the source text."

    return "\n".join(instructions)


# Main app logic
def main():
    
    glossary_data = [
      {"fr": "Personnes privées de titre de séjour ou sans titre de séjour", "nl": "Mensen zonder wettig verblijf", "en": "Undocumented person"},
      {"fr": "Personnes en séjour irrégulier", "nl": "Mensen zonder wettig verblijf", "en": "Undocumented person"},
      {"fr": "Personne privée de titre de séjour ou sans titre de séjour", "nl": "Persoon zonder wettig verblijf", "en": "Undocumented person"},
      {"fr": "Personnes usager.ère.s de drogues", "nl": "Drugsgebruikers", "en": "Drug user"},
      {"fr": "Personnes en situation de vulnérabilité", "nl": "Personen in maatschappelijk kwetsbare positie", "en": "People in vulnerable situations"},
      {"fr": "Personnes en situation de précarité", "nl": "Mensen in precaire situatie", "en": "People in precarious situations"},
      {"fr": "Personnes en situation de précarité", "nl": "Mensen in bestaansonzekerheid", "en": "People in precarious situations"},  
      {"fr": "Personnes en situation de pauvreté", "nl": "Mensen in armoede", "en": "People experiencing poverty"},
      {"fr": "Personnes en situation de pauvreté", "nl": "Mensen die in armoede leven", "en": "People living in poverty"},  
      {"fr": "Personnes refugiées", "nl": "Vluchtelingen", "en": "Refugee people"},
      {"fr": "Personnes en situation de migration", "nl": "Mensen met een migratieparcours", "en": "People in a migration situation"},
      {"fr": "Personne(s) sans abri et chez-soi", "nl": "Dak-en thuisloze persoon", "en": "Homeless people, person or people experiencing homelessness"},
      {"fr": "Sans-chez-soirisme", "nl": "dak- en thuisloosheid", "en": "Homelessness"},
      {"fr": "Aide médicale urgente", "nl": "Dringende medische hulp", "en": "Urgent medical aid"},
      {"fr": "AMU", "nl": "DMH", "en": "UMA"}, 
      {"fr": "Interruption volontaire de grossesse", "nl": "Vrijwillige zwangerschapsafbreking", "en": "Voluntary termination of pregnancy"},
      {"fr": "IVG", "nl": "VZA", "en": "VTP"},  
      {"fr": "Demandeur·euses de Protection Internationale (DPI)", "nl": "Verzoeker om Internationale Bescherming (VIB)", "en": "Applicant for international protection"},
      {"fr": "Les travailleur.euse.s du sexe (TDS)", "nl": "Sekswerkers", "en": "Sex workers"},
      {"fr": "Les Enfants et Jeunes en Situation de Rue (EJSR)", "nl": "Kinderen en jongeren die op straat leven", "en": "Children and Young People in Street Situations (CYPS)"},
      {"fr": "Mineur (Etranger) Non Accompagné (MENA)", "nl": "Niet begeleide minderjarige (vreemdeling) (NBMV)", "en": "(Foreign) unaccompanied minor"},
      {"fr": "Ayant(s)-droit", "nl": "De rechthebbende(n)", "en": "The right-holder(s)"},
      {"fr": "Partie prenante", "nl": "De belanghebbende(n)", "en": "Stakeholders"},
      {"fr": "Santé et droit sexuels et reproductifs (SDSR)", "nl": "Seksuele en reproductieve gezondheid en rechten (SRGR)", "en": "Sexual and Reproductive Health and Rigths (SRHR)"},
      {"fr": "Santé et droits en Migration", "nl": "Gezondheid, rechten en migratie", "en": "Health and Rights in Migration (HRM)"},
      {"fr": "Réduction de risques", "nl": "Harm Reduction", "en": "Harm Reduction"},
      {"fr": "Réduction de risques", "nl": "Risicobeperking", "en": "Harm Reduction"},
      {"fr": "Mesures de réduction des risques et des dommages", "nl": "Schade- en risicobeperkende maatregelen", "en": "Harm Reduction measures"},
      {"fr": "Salle de Consommation (à Moindre Risque)", "nl": "(risicobeperkende) gebruikersruimte", "en": "Supervised injection site"},
      {"fr": "Comptoir d'échange (de matériel de réduction des risques)", "nl": "Spuitenruil project", "en": "Syringe Service Programs"},
      {"fr": "Les programmes de drug checking", "nl": "Drugstest programma's", "en": "Drug checking programs"},
      {"fr": "Psycho-médico-social", "nl": "Psycho-medisch-sociaal", "en": ""},
      {"fr": "Soins de premier ligne", "nl": "Eerstelijnszorg", "en": "Primary care"},
      {"fr": "Soins de premier ligne", "nl": "Eerstelijnsgezondheidszorg", "en": "Primary care"},
      {"fr": "Santé mentale", "nl": "Geestelijke gezondheid", "en": "Mental health"},
      {"fr": "Soins de santé mentale", "nl": "Geestelijke gezondheidszorg", "en": "Mental health care"},
      {"fr": "SSM", "nl": "GGZ", "en": ""},  
      {"fr": "Services de Santé Mentale", "nl": "Centra voor Geestelijke Gezondheidszorg", "en": "Mental health services"},
      {"fr": "SSM", "nl": "CGZ", "en": ""},
      {"fr": "Problèmes de santé mentale", "nl": "Geestelijk gezondheidsproblemen", "en": "Mental health issues"},
      {"fr": "utilisateur de service", "nl": "zorggebruiker", "en": "Healthcare user"},
      {"fr": "Médecins du Monde", "nl": "Dokters van de Wereld", "en": "Doctors of the World"},
      {"fr": "CASO", "nl": "COZO", "en": ""}
    ]
        
    st.sidebar.title("Translation App")
    pass_word = st.sidebar.text_input('Enter the password:', type="password")
    if not pass_word:
        st.stop()
    if pass_word != PASSWORD:
        st.error('The password you entered is incorrect.')
        st.stop()

    
    tab1, tab2 = st.tabs(["Translate your Text", "Transform your text"])
    
    with tab1:
        
        st.header('Translate your text') 
        
        select_model = st.sidebar.radio("**Choose your tool**", ['Google Translate', 'Google Gemini', 'ChatGPT'])

        # Language selection
        target_lang = display_language_selection("target")
       
        
        if select_model == 'Google Translate':

            st.subheader('Google Translate')
            st.sidebar.write("Good for straightforward sentences and common phrases; it’s inexpensive and low-latency.")
            st.sidebar.write("**Drawbacks**: limited nuance and contextual understanding, difficulty with ambiguity and complex texts. Translations often mirror the source word-for-word.") 
            st.sidebar.info("For external use, add a refinement step or use a more advanced model.")
            st.sidebar.write("\n\n")
            
            uploaded_file = st.file_uploader("Upload file (PDF, PPTX, XLSX, DOCX)", type=['pdf', 'pptx', 'xlsx', 'docx'])
            file_text = ""
            
            if uploaded_file:              
                if uploaded_file.type == "application/pdf":
                    # file_text = read_pdf(uploaded_file)
                    uploaded_file_bytes = uploaded_file.getvalue()
                    file_text, page_images_bytes = process_uploaded_pdf_for_gemini(uploaded_file_bytes)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file_text = read_pptx(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    file_text = read_excel(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file_text = read_docx(uploaded_file)
                
            manual_text = display_text_input("basic_google")
            
            combined_text = file_text + "\n" + manual_text if file_text or manual_text else ""  
            
            if uploaded_file: 
                st.info(
                    "Text extracted from PDFs or PPTs may be dense or jumbled—lines can run together or appear out of order. "
                    "Use the feedback box to request layout improvements based on your instructions."
                )
                
            if 'translated_google_raw' not in st.session_state:
                st.session_state.translated_google_raw = ""
            
            if st.button('Translate', key="google_basic"):
                if combined_text:
                    
                    source_lang = detect_language(combined_text)

                    if source_lang[0] == target_lang:
                        st.error("Please select a target language other than the source language.")
                        st.stop()
                    
                    st.write(f"Detected language: {source_lang[1]}")
                    with st.spinner('translating...'):                     
                        translated_google = translate_text_with_v3(
                            text=combined_text, 
                            source_language_code=source_lang[0], 
                            target_language_code=target_lang,
                            project_id=project_id,
                            translate_client_obj=translate_client
                        )
                        st.session_state['translated_google_raw'] = translated_google
                        st.info("**Here is your translation:**")
                    pass
                                               
                else:
                    st.error('Please upload or paste a text to translate.')

            
            st.write(st.session_state['translated_google_raw'])

            if 'refined_google_translation' not in st.session_state:
                st.session_state.refined_google_translation = ""

            if 'refined_google_translation' not in st.session_state:
                st.session_state.refined_google_translation = ""
            
            if 'translated_google_raw' in st.session_state:
                st.write('**Refine with Gemini** ✨')
                guidelines = st.text_input("Your feedback and/or guidelines")
                        # Toon de download knop alleen als er tekst is om te downloaden (ruw of geformatteerd)
                if st.button('Start Refinement 🚀', key='basic_refine'):
                    with st.spinner('Refining... this can take a while'):
                        source_lang = detect_language(combined_text)
                        glossary_analysis = check_glossary(combined_text, source_lang, target_lang, glossary_data)
                        st.session_state.gloss_instruct = format_terminology_for_prompt(glossary_analysis, source_lang, target_lang)
                        st.write(st.session_state.gloss_instruct)
                        refined_gemini = improve_translation_with_gemini(combined_text, st.session_state['translated_google_raw'], guidelines, st.session_state.gloss_instruct, target_lang) 
                        st.session_state.refined_google_translation = refined_gemini
                        st.info("**GEMINI IMPROVED 🏋️**")
                    pass              
                
                st.write(st.session_state.refined_google_translation)
                                                                
            # Haal de teksten op uit session state die je in het document wilt zetten
            refined_text_for_download = st.session_state.get('refined_google_translation', '')
            raw_text_for_download = st.session_state.get('translated_google_raw', '')

            word_bytes = create_word_document(refined_text_for_download, raw_text_for_download) 

            if word_bytes is not None:
                st.sidebar.download_button(
                    label="Download latest translation (.docx)", # Pas label aan indien gewenst
                    data=word_bytes,
                    file_name="Jimmie_translation.docx", # Pas bestandsnaam aan indien gewenst
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
        
        elif select_model == 'Google Gemini':

            st.sidebar.write("Powered by **Gemini 2.5 Flash & Pro**, Google’s latest reasoning model. The quality will be similar to ChatGPT or to Google Basic tool after use of the Refinement button.") 
            st.sidebar.write("**Drawbacks:** results may vary—LLMs are nondeterministic and can hallucinate or misinterpret text. Refinement can take a while.")
            # st.sidebar.info("The prompt lets the model step back from the source to boost fluency, authenticity, and appeal.")
            
            st.write("Lower Temperature (~0.1 to 0.5): Recommended for more secure translations.")
            st.write("Higher Temperature (~0.6 to 1.0): Encourages more creative translations.")
            temp_choice = display_temperature_slider('gemini')

            uploaded_file = st.file_uploader("Upload file (PDF, PPTX, XLSX, DOCX)", type=['pdf', 'pptx', 'xlsx', 'docx'])
            file_text = ""
            
            if uploaded_file:
                if uploaded_file.type == "application/pdf":
                    file_text = read_pdf(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file_text = read_pptx(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    file_text = read_excel(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file_text = read_docx(uploaded_file)

            if uploaded_file: 
                st.info(
                    "Text extracted from PDFs or PPTs may be dense or jumbled—lines can run together or appear out of order. "
                    "You can use the feedback box to request layout improvements based on your instructions."
                )
            
            manual_text = display_text_input("gemini")

            combined_text = file_text + "\n" + manual_text if file_text or manual_text else None  

            
            # if 'translation_gemini' not in st.session_state:
            #     st.session_state.translation_gemini = ""
        
            translation_gemini = ""

            if 'translated_gemini_raw' not in st.session_state:
                st.session_state.translated_gemini_raw = ""
            
            if st.button('Translate', key='gemini_notbasic'):
            
                if combined_text:
                    source_lang = detect_language(combined_text)
                    st.write(f"Detected language: {source_lang[1]}")
                    with st.spinner('Neural network at work, be patient...'):
                        glossary_analysis = check_glossary(combined_text, source_lang[0], target_lang, glossary_data)
                        st.session_state.gloss_instruct = format_terminology_for_prompt(glossary_analysis, source_lang[0], target_lang)
                        st.session_state.translated_gemini_raw = translate_text_vertexai_prompting(combined_text, source_lang, target_lang, st.session_state.gloss_instruct, temp_choice, "gemini-2.5-flash-preview-05-20")
                        st.info("**GEMINI TRANSLATED**")
                else:
                    st.error('Please upload or paste a text to translate.')
            
            st.write(st.session_state.translated_gemini_raw)
            
            if 'refined_translation' not in st.session_state:
                st.session_state.refined_translation = ""
            
            if 'translated_gemini_raw' in st.session_state:
                st.write('**Refine with Gemini** ✨')
                guidelines = st.text_input("Your feedback and/or guidelines")
                        
                if st.button('Start Refinement 🚀', key='gemini_refine'):
                    with st.spinner('👉 Thinking hard... this can take a while.... '):
                        refined_gemini = improve_translation_with_gemini(combined_text, st.session_state['translated_gemini_raw'], guidelines, st.session_state.gloss_instruct, target_lang) 
                        st.session_state.refined_translation = refined_gemini
                        st.info("**GEMINI IMPROVED 🏋️**")
                    pass              
                
                st.write(st.session_state.refined_translation)
                                                                
            # Haal de teksten op uit session state die je in het document wilt zetten
            refined_text_for_download = st.session_state.get('refined_translation', '')
            raw_text_for_download = st.session_state.get('translated_gemini_raw', '')

            word_bytes = create_word_document(refined_text_for_download, raw_text_for_download) 

            if word_bytes is not None:
                st.sidebar.download_button(
                    label="Download latest translation (.docx)", # Pas label aan indien gewenst
                    data=word_bytes,
                    file_name="Jimmie_translation.docx", # Pas bestandsnaam aan indien gewenst
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

        elif select_model == 'ChatGPT':

            st.sidebar.write("Powered by **GPT 4.1 & o3**, OpenAI's latest models. The quality will be similar to Google Gemini but at a higher cost. Can be used to challenge the Gemini translation.") 
            st.sidebar.write("**Drawbacks:** results may vary—LLMs are nondeterministic and can hallucinate or misinterpret text.")
            # st.sidebar.info("The prompt lets the model step back from the source to boost fluency, authenticity, and appeal.")
            
            
            st.write("Lower Temperature (~0.1 to 0.5): Recommended for more secure translations.")
            st.write("Higher Temperature (~0.6 to 1.0): Encourages more creative translations.")
            temp_choice = display_temperature_slider('gpt')
            
            uploaded_file = st.file_uploader("Upload file (PDF, PPTX, XLSX, DOCX)", type=['pdf', 'pptx', 'xlsx', 'docx'])
            file_text = ""
            
            if uploaded_file:
                if uploaded_file.type == "application/pdf":
                    file_text = read_pdf(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file_text = read_pptx(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    file_text = read_excel(uploaded_file)
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file_text = read_docx(uploaded_file)

            if uploaded_file: 
                st.info(
                    "Text extracted from PDFs or PPTs may be dense or jumbled—lines can run together or appear out of order. "
                    "You can use the feedback box to request layout improvements based on your instructions."
                )
            
            manual_text = display_text_input("openai")

            combined_text = file_text + "\n" + manual_text if file_text or manual_text else None      
           
            if 'translated_gpt_raw' not in st.session_state:
                st.session_state.translated_gpt_raw = ""
            
            if st.button('Translate', key='openai_o3'):
                
                if combined_text:
                    source_lang = detect_language(combined_text)

                    if source_lang[0] == target_lang:
                        st.error("Please select a target language other than the source language.")
                        st.stop()  # stopt de rest van het script hier
                        
                    st.write(f"Detected language: {source_lang[1]}")
                    with st.spinner('👉 Thinking hard... Please wait a moment.'):
                        basic_google = translate_text_with_v3(
                            text=combined_text, 
                            source_language_code=source_lang[0], 
                            target_language_code=target_lang,
                            project_id=project_id,
                            translate_client_obj=translate_client
                        )
                        basic_openai = get_openai_translation(combined_text, source_lang[0], target_lang, temp_choice, model="gpt-4.1")
                        glossary_analysis = check_glossary(combined_text, source_lang[0], target_lang, glossary_data)
                        st.session_state.gloss_instruct = format_terminology_for_prompt(glossary_analysis, source_lang[0], target_lang)
                        st.session_state.translated_gpt_raw = get_openai_refinement(basic_google, basic_openai, target_lang, st.session_state.gloss_instruct)
                        st.info("**ChatGPT TRANSLATED**")
                else:
                    st.error('Please upload or paste a text to translate.')
                
                
                st.write(st.session_state.translated_gpt_raw)
            
            if 'refined_gpt_translation' not in st.session_state:
                st.session_state.refined_gpt_translation = ""
            
            if 'translated_gpt_raw' in st.session_state:
                st.write('**Refine with o3** ✨')
                guidelines = st.text_input("Your feedback and/or guidelines")
                        # Toon de download knop alleen als er tekst is om te downloaden (ruw of geformatteerd)
                if st.button('Start Refinement 🚀', key="openai_refine"):
                    with st.spinner('Thinking hardddd... '):
                        
                        refined_gpt = improve_translation_with_gpt(combined_text, st.session_state.translated_gpt_raw, guidelines, target_lang, st.session_state.gloss_instruct) 
                        st.session_state.refined_gpt_translation = refined_gpt
                        st.info("**o3 IMPROVED 🏋️**")
                    pass              
                
                st.write(st.session_state.refined_gpt_translation)
                                                                
            # Haal de teksten op uit session state die je in het document wilt zetten
            refined_text_for_download = st.session_state.get('refined_gpt_translation', '')
            raw_text_for_download = st.session_state.get('translated_gpt_raw', '')

            word_bytes = create_word_document(refined_text_for_download, raw_text_for_download) 

            if word_bytes is not None:
                st.sidebar.download_button(
                    label="Download latest translation (.docx)", # Pas label aan indien gewenst
                    data=word_bytes,
                    file_name="Jimmie_translation.docx", # Pas bestandsnaam aan indien gewenst
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
            
        if st.sidebar.button("Clear all translations"):
            st.session_state.refined_translation = ""
            st.session_state.translated_google_raw = ""
            st.session_state.refined_google_translation = ""
            st.session_state.refined_gpt_translation = ""
            st.session_state.translated_gpt_raw = ""
            st.session_state.translated_gemini_raw = ""
            st.rerun()     

    with tab2:
        st.write("🚧 under construction 🚧")
    
    # with tab3:
        
    #     select_model = st.sidebar.radio('**Select your MODEL**', ['gpt-4.1', 'MISTRAL large'])
    #     tool_choice = st.sidebar.radio("**Choose your tool:**", ['Single Agent', 'Multi-Agent'])
    #     st.sidebar.write("*The multi-agent system is likely to produce better results, albeit with a higher footprint and longer runtime.*")
    #     st.sidebar.write("*Making smart use of the feedback mechanisms can yield great results. Give it a try.*")
    #     st.sidebar.write("*The third tool is under construction.*")
        
    #     if tool_choice == 'Single Agent':
    #         translate_with_enhancement(select_model)
    #     elif tool_choice == 'Multi-Agent':
    #         multiagent_translation(select_model)
               
    #     manage_central_file()

if __name__ == "__main__":
    main()


